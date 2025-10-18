"""Convert HTML slides to PPTX presentation (using rendered images)."""

from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches
from .utils import get_logger, ensure_directory
from .html_to_image import convert_html_files_to_images

logger = get_logger(__name__)


class HTMLToPPTXConverter:
    """Converts HTML slides to PPTX presentation using rendered images."""

    def __init__(self):
        """Initialize converter."""
        logger.info("HTMLToPPTXConverter initialized (image-based)")

    def convert_html_files_to_pptx(
        self,
        html_files: List[Path],
        output_path: Path,
        image_dir: Optional[Path] = None
    ) -> Path:
        """Convert HTML files to PPTX by rendering them as images.

        Args:
            html_files: List of HTML file paths (in order)
            output_path: Path to save PPTX file
            image_dir: Directory to save rendered images (default: same as output_path)

        Returns:
            Path to generated PPTX file
        """
        logger.info(f"Converting {len(html_files)} HTML files to PPTX (image-based)")

        # Determine image directory
        if image_dir is None:
            image_dir = output_path.parent / "slide_images"
        ensure_directory(image_dir)

        # Step 1: Convert HTML files to images using Playwright
        logger.info("Step 1: Rendering HTML files to PNG images...")
        image_paths = convert_html_files_to_images(
            html_files=html_files,
            output_dir=image_dir,
            width=1920,   # Higher resolution for better quality
            height=1080,  # 16:9 aspect ratio
            headless=True
        )

        if not image_paths:
            raise ValueError("No images were generated from HTML files")

        logger.info(f"Generated {len(image_paths)} slide images")

        # Step 2: Create PPTX and insert images as slides
        logger.info("Step 2: Creating PPTX presentation...")
        prs = Presentation()
        
        # Set slide dimensions to 16:9 (10 x 5.625 inches)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        # Use blank layout for full-slide images
        blank_layout = prs.slide_layouts[6]  # Blank layout

        for i, image_path in enumerate(image_paths):
            logger.debug(f"Adding slide {i+1}/{len(image_paths)}: {image_path.name}")

            # Add blank slide
            slide = prs.slides.add_slide(blank_layout)

            # Add image to fill entire slide
            # Position: (0, 0) - top-left corner
            # Size: full slide dimensions
            left = Inches(0)
            top = Inches(0)
            width = prs.slide_width
            height = prs.slide_height

            try:
                slide.shapes.add_picture(
                    str(image_path),
                    left,
                    top,
                    width=width,
                    height=height
                )
                logger.info(f"✓ Added slide {i+1}: {image_path.name}")

            except Exception as e:
                logger.error(f"Failed to add image {image_path}: {e}")
                raise

        # Step 3: Save presentation
        ensure_directory(output_path.parent)
        prs.save(str(output_path))

        logger.info(f"✓ PPTX saved to: {output_path}")
        logger.info(f"   • Slides: {len(image_paths)}")
        logger.info(f"   • Images: {image_dir}")

        return output_path


# Convenience function
def convert_html_to_pptx(
    html_files: List[Path],
    output_path: Path,
    image_dir: Optional[Path] = None
) -> Path:
    """Convert HTML slides to PPTX.

    Args:
        html_files: List of HTML file paths
        output_path: Output PPTX path
        image_dir: Directory for rendered images

    Returns:
        Path to generated PPTX file
    """
    converter = HTMLToPPTXConverter()
    return converter.convert_html_files_to_pptx(
        html_files,
        output_path,
        image_dir
    )
