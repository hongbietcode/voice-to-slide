"""Slide building module using python-pptx."""

from pathlib import Path
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from .utils import get_logger, ensure_directory

logger = get_logger(__name__)

class SlideBuilder:
    """Builds PowerPoint presentations using python-pptx."""

    # Default color scheme
    TITLE_COLOR = RGBColor(44, 62, 80)      # Dark blue-gray
    TEXT_COLOR = RGBColor(52, 73, 94)       # Slightly lighter blue-gray
    ACCENT_COLOR = RGBColor(41, 128, 185)   # Blue accent
    BACKGROUND_COLOR = RGBColor(255, 255, 255)  # White

    def __init__(self, template_path: Optional[Path | str] = None):
        """Initialize the slide builder.

        Args:
            template_path: Optional path to PPTX template file
        """
        if template_path and Path(template_path).exists():
            self.prs = Presentation(str(template_path))
            logger.info(f"Loaded template from {template_path}")
        else:
            self.prs = Presentation()
            logger.info("Created new presentation")

        # Set presentation dimensions (16:9 widescreen)
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        """Add a title slide to the presentation.

        Args:
            title: Main title text
            subtitle: Optional subtitle text
        """
        logger.info(f"Adding title slide: {title}")

        # Use the first slide layout (typically title slide)
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.TITLE_COLOR

        # Set subtitle if available
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = self.TEXT_COLOR

    def add_content_slide(
        self,
        title: str,
        bullet_points: List[str],
        image_path: Optional[Path | str] = None
    ) -> None:
        """Add a content slide with title, bullet points, and optional image.

        Args:
            title: Slide title
            bullet_points: List of bullet points
            image_path: Optional path to image file
        """
        logger.info(f"Adding content slide: {title}")

        # Create blank slide
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)

        # Add title
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.8)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.TITLE_COLOR

        # Determine content area based on whether there's an image
        if image_path and Path(image_path).exists():
            # Split slide: text on left, image on right
            content_left = Inches(0.5)
            content_top = Inches(1.3)
            content_width = Inches(4.5)
            content_height = Inches(3.8)

            image_left = Inches(5.2)
            image_top = Inches(1.3)
            image_width = Inches(4.3)
            image_height = Inches(3.8)

            # Add image
            try:
                slide.shapes.add_picture(
                    str(image_path),
                    image_left,
                    image_top,
                    width=image_width,
                    height=image_height
                )
                logger.info(f"Added image: {image_path}")
            except Exception as e:
                logger.error(f"Failed to add image {image_path}: {e}")
        else:
            # Full width for text
            content_left = Inches(0.5)
            content_top = Inches(1.3)
            content_width = Inches(9)
            content_height = Inches(3.8)

        # Add bullet points
        text_box = slide.shapes.add_textbox(
            content_left,
            content_top,
            content_width,
            content_height
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_before = Pt(12) if i > 0 else Pt(0)

    def add_section_slide(self, section_title: str) -> None:
        """Add a section divider slide.

        Args:
            section_title: Section title text
        """
        logger.info(f"Adding section slide: {section_title}")

        # Use title-only layout
        slide_layout = self.prs.slide_layouts[5]  # Title only
        slide = self.prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = section_title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.ACCENT_COLOR
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def build_presentation(
        self,
        content: Dict[str, Any],
        image_paths: List[Optional[Path]] = None
    ) -> None:
        """Build complete presentation from content structure.

        Args:
            content: Content dictionary with 'title' and 'slides' keys
            image_paths: Optional list of image paths for each slide
        """
        logger.info("Building complete presentation")

        # Add title slide
        presentation_title = content.get("title", "Presentation")
        self.add_title_slide(presentation_title)

        # Add content slides
        slides = content.get("slides", [])
        if image_paths is None:
            image_paths = [None] * len(slides)

        for i, slide_data in enumerate(slides):
            title = slide_data.get("title", f"Slide {i+1}")
            bullet_points = slide_data.get("bullet_points", [])
            image_path = image_paths[i] if i < len(image_paths) else None

            self.add_content_slide(title, bullet_points, image_path)

        logger.info(f"Presentation built with {len(self.prs.slides)} slides")

    def save(self, output_path: Path | str) -> Path:
        """Save the presentation to a file.

        Args:
            output_path: Path to save the PPTX file

        Returns:
            Path to the saved file
        """
        output_path = Path(output_path)
        ensure_directory(output_path.parent)

        self.prs.save(str(output_path))
        logger.info(f"Presentation saved to {output_path}")

        return output_path

    @staticmethod
    def create_presentation(
        content: Dict[str, Any],
        output_path: Path | str,
        image_paths: List[Optional[Path]] = None,
        template_path: Optional[Path | str] = None
    ) -> Path:
        """Convenience method to create and save a presentation.

        Args:
            content: Content dictionary with 'title' and 'slides' keys
            output_path: Path to save the PPTX file
            image_paths: Optional list of image paths for each slide
            template_path: Optional path to PPTX template file

        Returns:
            Path to the saved file
        """
        builder = SlideBuilder(template_path)
        builder.build_presentation(content, image_paths)
        return builder.save(output_path)
