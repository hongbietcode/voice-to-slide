"""Convert HTML slides to PPTX presentation (using rendered images)."""

from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches
from .utils import get_logger, ensure_directory
from .html_to_image import convert_html_files_to_images

logger = get_logger(__name__)


class HTMLToPPTXConverter:
    """Converts HTML slides to PPTX presentation using rendered images."""

    def __init__(self):
        """Initialize converter."""
        logger.info("HTMLToPPTXConverter initialized (image-based)")

    def parse_html_file(self, html_path: Path) -> Dict[str, Any]:
        """Parse an HTML file and extract slide content.

        Args:
            html_path: Path to HTML file

        Returns:
            Dictionary with title, bullet_points, image_src
        """
        try:
            with open(html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            parser = HTMLSlideParser()
            parser.feed(html_content)

            return {
                "title": parser.title,
                "bullet_points": parser.bullet_points,
                "image_src": parser.image_src
            }

        except Exception as e:
            logger.error(f"Failed to parse HTML file {html_path}: {e}")
            raise

    def _parse_color_from_css(self, css_content: str, property_name: str) -> Optional[RGBColor]:
        """Extract color from CSS content.

        Args:
            css_content: CSS string
            property_name: CSS property name (e.g., 'color', 'background')

        Returns:
            RGBColor if found, None otherwise
        """
        # Look for hex colors like #2C3E50
        pattern = rf'{property_name}:\s*#([0-9A-Fa-f]{{6}})'
        match = re.search(pattern, css_content)

        if match:
            hex_color = match.group(1)
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)

        # Look for rgb colors like rgb(44, 62, 80)
        pattern = rf'{property_name}:\s*rgb\((\d+),\s*(\d+),\s*(\d+)\)'
        match = re.search(pattern, css_content)

        if match:
            r, g, b = int(match.group(1)), int(match.group(2)), int(match.group(3))
            return RGBColor(r, g, b)

        return None

    def extract_theme_from_html(self, html_path: Path) -> Dict[str, RGBColor]:
        """Extract theme colors from HTML file CSS.

        Args:
            html_path: Path to HTML file

        Returns:
            Dictionary of color names to RGBColor objects
        """
        try:
            with open(html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            # Extract style tag content
            style_match = re.search(r'<style>(.*?)</style>', html_content, re.DOTALL)
            if not style_match:
                logger.warning(f"No <style> tag found in {html_path}, using defaults")
                return self.DEFAULT_COLORS

            css_content = style_match.group(1)

            # Extract colors
            theme_colors = {}

            # Try to find colors used in the CSS
            if primary := self._parse_color_from_css(css_content, 'color'):
                theme_colors["text"] = primary

            # Look for .slide-title color
            title_pattern = r'\.slide-title\s*\{[^}]*color:\s*#([0-9A-Fa-f]{6})'
            if match := re.search(title_pattern, css_content):
                hex_color = match.group(1)
                r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
                theme_colors["primary"] = RGBColor(r, g, b)

            # Look for accent color (in li::before)
            accent_pattern = r'li::before\s*\{[^}]*color:\s*#([0-9A-Fa-f]{6})'
            if match := re.search(accent_pattern, css_content):
                hex_color = match.group(1)
                r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
                theme_colors["accent"] = RGBColor(r, g, b)

            # Merge with defaults
            return {**self.DEFAULT_COLORS, **theme_colors}

        except Exception as e:
            logger.warning(f"Failed to extract theme from {html_path}: {e}, using defaults")
            return self.DEFAULT_COLORS

    def convert_html_files_to_pptx(
        self,
        html_files: List[Path],
        output_path: Path,
        extract_theme: bool = True
    ) -> Path:
        """Convert a list of HTML files to a PPTX presentation.

        Args:
            html_files: List of HTML file paths (in order)
            output_path: Path to save PPTX file
            extract_theme: Whether to extract theme colors from HTML

        Returns:
            Path to generated PPTX file
        """
        logger.info(f"Converting {len(html_files)} HTML files to PPTX")

        # Extract theme from first HTML file if requested
        if extract_theme and html_files:
            self.colors = self.extract_theme_from_html(html_files[0])
            logger.info("Extracted theme colors from HTML")

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        # Process each HTML file
        for i, html_file in enumerate(html_files):
            logger.info(f"Processing slide {i+1}/{len(html_files)}: {html_file.name}")

            try:
                slide_data = self.parse_html_file(html_file)

                if i == 0:
                    # Title slide
                    self._add_title_slide(prs, slide_data)
                else:
                    # Content slide
                    self._add_content_slide(prs, slide_data)

            except Exception as e:
                logger.error(f"Failed to process {html_file}: {e}")
                # Continue with other slides

        # Save presentation
        ensure_directory(output_path.parent)
        prs.save(str(output_path))

        logger.info(f"PPTX saved to: {output_path}")
        return output_path

    def _add_title_slide(self, prs: Presentation, slide_data: Dict[str, Any]) -> None:
        """Add title slide to presentation.

        Args:
            prs: Presentation object
            slide_data: Parsed slide data
        """
        # Use title layout
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = slide_data.get("title", "")
        title_frame = title.text_frame
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors.get("primary", self.DEFAULT_COLORS["primary"])
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        logger.debug(f"Added title slide: {slide_data.get('title', '')}")

    def _add_content_slide(self, prs: Presentation, slide_data: Dict[str, Any]) -> None:
        """Add content slide to presentation.

        Args:
            prs: Presentation object
            slide_data: Parsed slide data
        """
        # Use blank layout for custom positioning
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add title
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.8)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get("title", "")
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors.get("primary", self.DEFAULT_COLORS["primary"])

        # Check if slide has image
        image_src = slide_data.get("image_src")
        has_image = image_src and Path(image_src).exists()

        if has_image:
            # Two-column layout
            content_left = Inches(0.5)
            content_top = Inches(1.3)
            content_width = Inches(4.5)
            content_height = Inches(3.8)

            image_left = Inches(5.2)
            image_top = Inches(1.3)
            image_width = Inches(4.3)
            image_height = Inches(3.8)

            # Add image
            try:
                slide.shapes.add_picture(
                    image_src,
                    image_left,
                    image_top,
                    width=image_width,
                    height=image_height
                )
                logger.debug(f"Added image: {image_src}")
            except Exception as e:
                logger.error(f"Failed to add image {image_src}: {e}")
                # Fall back to full-width text
                has_image = False

        if not has_image:
            # Full-width text
            content_left = Inches(0.5)
            content_top = Inches(1.3)
            content_width = Inches(9)
            content_height = Inches(3.8)

        # Add bullet points
        text_box = slide.shapes.add_textbox(
            content_left,
            content_top,
            content_width,
            content_height
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        bullet_points = slide_data.get("bullet_points", [])
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors.get("text", self.DEFAULT_COLORS["text"])
            p.space_before = Pt(12) if i > 0 else Pt(0)

        logger.debug(f"Added content slide with {len(bullet_points)} bullets")


# Convenience function
def convert_html_to_pptx(
    html_files: List[Path],
    output_path: Path,
    theme_colors: Optional[Dict[str, RGBColor]] = None,
    extract_theme: bool = True
) -> Path:
    """Convert HTML slides to PPTX.

    Args:
        html_files: List of HTML file paths
        output_path: Output PPTX path
        theme_colors: Optional theme color overrides
        extract_theme: Extract theme from HTML

    Returns:
        Path to generated PPTX file
    """
    converter = HTMLToPPTXConverter(theme_colors=theme_colors)
    return converter.convert_html_files_to_pptx(html_files, output_path, extract_theme)
